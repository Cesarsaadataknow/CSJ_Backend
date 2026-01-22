from datetime import datetime
from pytz import timezone
from docx import Document
from pptx import Presentation
import pandas as pd
import timeit
import uuid
import tiktoken
import pdb
import io


def genereta_id() -> str: 
    now = datetime.now()
    str_now = now.strftime("%Y%m%d")
    uuid_id = str(uuid.uuid4())
    chat_id = f'{str_now}-{uuid_id}'
    return chat_id

def current_colombian_time() -> str:
    current_time = datetime.now(timezone('America/Bogota')).strftime('%Y-%m-%d %H:%M:%S')
    return current_time

def timeit_decorator(func):
    def wrapper(*args, **kwargs):
        start_time = timeit.default_timer()
        result = func(*args, **kwargs)
        end_time = timeit.default_timer()
        elapsed_time = end_time - start_time
        return result, elapsed_time
    return wrapper

def count_tokens(texts=None, model_reference="cl100k_base"):   
    if texts:
        encoding = tiktoken.get_encoding(model_reference)
        count = encoding.encode(texts)
        return count
    
def format_conversation_data(documents):
    # Suponiendo que todos son de la misma conversación
    if not documents:
        return None

    conversation_id = documents[0].get('conversation_id')
    messages = []
    cutoff_time = None  # Fecha de actualización del documento con flag_modifier
    for doc in documents:
        user_msg = doc.get("user_message", {})
        files = doc.get("pdf_text") 
        ai_msg = doc.get("ai_message", {})
        # Convertimos la fecha de creación del mensaje del usuario a datetime
        created_at_str = user_msg.get("created_at")
        if created_at_str:
            created_at = datetime.fromisoformat(created_at_str)
        else:
            # Si no se encuentra la fecha, se asume que se procesa el documento
            created_at = None

        # Si ya se estableció un cutoff_time y el documento tiene una fecha de creación,
        # omitimos el documento si su fecha de creación es anterior a cutoff_time.
        if cutoff_time and created_at and created_at < cutoff_time:
            continue

        # Mensaje de usuario
        messages.append({
            "id": doc.get("id"),
            "role": "user",
            "content": user_msg.get("content"),
            "created_at": user_msg.get("created_at"),
            "files": doc.get("files_in_message")  
        })

        # Definimos el contenido del mensaje de la IA, considerando la condición de rate
        if doc.get('rate') == 2:
            ai_content = "Mensaje cancelado por el usuario"
        else:
            ai_content = ai_msg.get("content")

        # Mensaje de la IA
        messages.append({
            "id": doc.get("id"),
            "role": "assistant",
            "content": ai_content,
            "created_at": ai_msg.get("created_at"),
            "rate": doc.get("rate")
        })

        # Si se detecta flag_modifier en True y aún no se ha establecido cutoff_time,
        # se guarda la fecha de actualización (updated_at) para usarla como fecha de corte.
        if doc.get("flag_modifier") is True and cutoff_time is None:
            updated_at_str = doc.get("updated_at")
            if updated_at_str:
                cutoff_time = datetime.fromisoformat(updated_at_str)

    # Ordenar los mensajes por fecha de creación
    messages.sort(key=lambda msg: datetime.fromisoformat(msg["created_at"]))

    return {
        "conversation_id": conversation_id,
        "conversation_name": documents[0].get("conversation_name"),
        "user_id": documents[0].get("user_id"),
        "messages": messages
    }

def extract_text_content(content: bytes) -> str:
    # Asumimos que el contenido está en UTF-8
    return content.decode("utf-8")

def extract_word_content(content: bytes) -> str:
    """
    Extrae el texto de un archivo Word (DOCX) a partir de un objeto de bytes.

    Args:
        content (bytes): Contenido del archivo Word en formato bytes.

    Returns:
        str: Texto extraído del archivo Word.
    """
    # Convertir los bytes en un stream de memoria
    stream = io.BytesIO(content)
    # Abrir el documento usando python-docx
    document = Document(stream)
    
    # Extraer el texto de cada párrafo del documento
    full_text = []
    for paragraph in document.paragraphs:
        full_text.append(paragraph.text)
    
    # Unir los párrafos en un único string, separándolos por saltos de línea
    return "\n".join(full_text)

def extract_excel_content(content: bytes) -> str:
    """
    Lee el contenido de un archivo Excel (en formato bytes) y lo convierte a JSONL,
    donde cada línea representa un registro del Excel.

    Args:
        content (bytes): Contenido del archivo Excel en formato bytes.

    Returns:
        str: Un string en formato JSONL con los registros del Excel.
    """
    try:
        # Crear un objeto BytesIO a partir de los bytes del archivo Excel
        excel_io = io.BytesIO(content)
        # Leer el contenido del Excel en un DataFrame de pandas
        df = pd.read_excel(excel_io)
    except Exception as e:
        raise ValueError(f"Error al leer el archivo Excel: {str(e)}")
    
    # Convertir el DataFrame a JSON Lines (cada registro en una línea)
    jsonl_string = df.to_json(orient="records", lines=True, force_ascii=False)
    return jsonl_string

def extract_pptx_content(content: bytes) -> str:
    """
    Extrae el texto de un archivo PPTX a partir de un objeto de bytes.
    """
    try:
        stream = io.BytesIO(content)
        presentation = Presentation(stream)

        # Recorremos cada diapositiva y cada shape que tenga texto
        slides_text = []
        for slide in presentation.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text)
            # Unimos el texto de la diapositiva con saltos de línea
            slides_text.append("\n".join(slide_texts))
        # Unimos todas las diapositivas
        full_text = "\n".join(slides_text)
        return full_text

    except Exception as e:
        raise ValueError(f"Error al leer el archivo PPTX: {str(e)}")

    
